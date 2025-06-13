from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

"""
PPTX Presentation Generation Module

This module provides functionality to dynamically generate PowerPoint presentations (.pptx) from structured data.

The `PPTXService` class contains a static method to create presentations using the `python-pptx` library.

Classes:
    - PPTXService: Provides a utility method for creating presentations with slide titles and content.

Functions:
    - create_presentation(topic, slides, file_path=None): Generates and saves a .pptx file using the provided topic and slide data.

Slide Data Format:
    - The `slides` parameter should be a list where each item is either:
        1. A string (used as slide content under a default title "Slide")
        2. A dictionary with:
            - "title": (str) Slide title
            - "bullet_points": (list of str) Content as bullet points

Workflow:
    1. Create a new PowerPoint presentation using the default layout.
    2. For each slide in the input list:
        - Add a slide with a title and content (text or bullet points).
    3. Save the final presentation to the specified path.

Usage Example:
    slides = [
        "Welcome to the presentation!",
        {"title": "Introduction", "bullet_points": ["Overview", "Objectives"]},
        {"title": "Conclusion", "bullet_points": ["Summary", "Q&A"]}
    ]
    file_path = PPTXService.create_presentation("My Topic", slides)

Parameters:
    - topic (str): The title/topic used to name the file (if `file_path` is not provided).
    - slides (list): List of slide contents (string or dictionary format).
    - file_path (str, optional): Custom path to save the presentation.

Returns:
    - str: Full path to the saved `.pptx` file.

Note:
    - Default file path is `/tmp/{topic}.pptx` if not specified.
    - Uses layout index 1 (Title and Content) from the PowerPoint template.
"""


class PPTXService:
    @staticmethod
    def create_presentation(topic, slides, file_path=None):
        if file_path is None:
            file_path = f"/tmp/{topic.replace(' ', '_')}.pptx"

        prs = Presentation()
        bullet_slide_layout = prs.slide_layouts[1]

        for slide_data in slides:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default content

            if isinstance(slide_data, str):
                title_shape.text = "Slide"
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = slide_data
                run.font.size = Pt(24)
                run.font.bold = True

            elif isinstance(slide_data, dict):
                title_shape.text = slide_data.get("title", "Untitled")

                # Optional Subtitle
                if "subtitle" in slide_data:
                    p = text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = slide_data["subtitle"]
                    run.font.size = Pt(18)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(100, 100, 100)

                # Bullet Points
                for point in slide_data.get("bullet_points", []):
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(20)

                image_path = slide_data.get("image_path")
                if image_path and os.path.exists(image_path):
                    try:
                        slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4.5))
                    except Exception as e:
                        print(f"Could not insert image: {e}")

        prs.save(file_path)
        return file_path
